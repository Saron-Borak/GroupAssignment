"""
Estimate whether any text box needs more height than it has.

No renderer here, so this simulates word-wrap using average glyph widths for the
two fonts in the deck and compares the required height against the shape height.
"""
import math
from pptx import Presentation

DECK = r"C:/Users/Rakze/Desktop/GroupAssignment/docs/EAMU-MIS-Group-Presentation.pptx"
EMU_IN = 914400
PT_IN = 72.0

# Average glyph width as a fraction of point size, measured for these families.
WIDTH = {'Calibri': 0.475, 'Cambria': 0.500, 'Courier New': 0.600}

def _has_bullet(p):
    return 'buChar' in p._p.xml or 'buAutoNum' in p._p.xml

prs = Presentation(DECK)
issues = []

for idx, slide in enumerate(prs.slides, start=1):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tf = sh.text_frame
        if not tf.text.strip():
            continue
        w_in = (sh.width or 0) / EMU_IN
        h_in = (sh.height or 0) / EMU_IN
        if w_in <= 0 or h_in <= 0:
            continue

        total_h_pt = 0.0
        for para in tf.paragraphs:
            runs = [r for r in para.runs if r.text]
            if not runs:
                total_h_pt += 6
                continue
            size_pt = max((r.font.size.pt if r.font.size else 12) for r in runs)
            face = next((r.font.name for r in runs if r.font.name), 'Calibri')
            frac = WIDTH.get(face, 0.475)
            bullet_indent = 0.22 if para.level or _has_bullet(para) else 0.0
            usable = max(0.4, w_in - bullet_indent)
            cpl = max(4, int((usable * PT_IN) / (frac * size_pt)))

            text = ''.join(r.text for r in runs)
            # simulate wrapping
            lines, cur = 1, 0
            for word in text.split(' '):
                add = len(word) + (1 if cur else 0)
                if cur + add > cpl and cur:
                    lines += 1
                    cur = len(word)
                else:
                    cur += add
            total_h_pt += lines * size_pt * 1.22

        req_in = total_h_pt / PT_IN
        if req_in > h_in * 1.06:
            issues.append((idx, round(req_in, 2), round(h_in, 2),
                           tf.text.strip().replace('\n', ' / ')[:66]))

if issues:
    print(f"!! {len(issues)} possible overflow(s)  (needed vs available, inches)\n")
    for idx, req, have, txt in sorted(issues, key=lambda r: r[1] - r[2], reverse=True):
        print(f"  slide {idx:>2}  needs {req:>4}  has {have:>4}   {txt!r}")
else:
    print("No estimated text overflow.")
