"""
Geometry QA for the deck.

No renderer is available on this machine, so instead of eyeballing images this
checks the two defects that would actually be visible: shapes crossing the slide
edge, and text boxes overlapping each other.
"""
from pptx import Presentation
from pptx.util import Emu

DECK = r"C:/Users/Rakze/Desktop/GroupAssignment/docs/EAMU-MIS-Group-Presentation.pptx"
prs = Presentation(DECK)

SW, SH = prs.slide_width, prs.slide_height
EMU_IN = 914400
MARGIN = int(0.4 * EMU_IN)          # tolerated minimum edge gap

def inches(v):
    return round(v / EMU_IN, 2)

print(f"slide size: {inches(SW)} x {inches(SH)} in")
print(f"slides: {len(prs.slides)}\n")

problems = []

for idx, slide in enumerate(prs.slides, start=1):
    boxes = []
    for sh in slide.shapes:
        if sh.left is None or sh.top is None:
            continue
        l, t = sh.left, sh.top
        r, b = l + (sh.width or 0), t + (sh.height or 0)

        # 1. Off-slide or too close to the edge
        if l < 0 or t < 0 or r > SW or b > SH:
            problems.append(
                f"slide {idx}: OFF-SLIDE {sh.shape_type} at "
                f"({inches(l)},{inches(t)})-({inches(r)},{inches(b)})"
            )
        elif l < MARGIN or t < MARGIN or (SW - r) < MARGIN or (SH - b) < MARGIN:
            # The page-number label sits low by design.
            txt = (sh.text_frame.text[:20] if sh.has_text_frame else "")
            if not txt.strip().isdigit():
                problems.append(
                    f"slide {idx}: TIGHT MARGIN {inches(min(l, t, SW-r, SH-b))}in  "
                    f"text={txt!r}"
                )

        if sh.has_text_frame and sh.text_frame.text.strip():
            boxes.append((l, t, r, b, sh.text_frame.text.strip()[:34]))

    # 2. Text boxes overlapping other text boxes
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            a, bx = boxes[i], boxes[j]
            ox = min(a[2], bx[2]) - max(a[0], bx[0])
            oy = min(a[3], bx[3]) - max(a[1], bx[1])
            if ox > int(0.06 * EMU_IN) and oy > int(0.06 * EMU_IN):
                problems.append(
                    f"slide {idx}: TEXT OVERLAP {inches(ox)}x{inches(oy)}in  "
                    f"{a[4]!r} <> {bx[4]!r}"
                )

if problems:
    print(f"!! {len(problems)} issue(s)\n")
    for p in problems:
        print("  " + p)
else:
    print("No off-slide shapes, tight margins or text overlaps found.")
